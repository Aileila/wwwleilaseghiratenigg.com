from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# LSN tokens
BLUE = RGBColor(0x33, 0x33, 0xCC)
INK = RGBColor(0x2A, 0x2A, 0x34)
CREAM = RGBColor(0xED, 0xE9, 0xE3)
TAUPE = RGBColor(0x94, 0x8F, 0x8A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
# 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, text, color=INK, top=0.6, size=46):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.8), Inches(1.3))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Noto Serif"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = color
    return box


def add_subtitle(slide, text, color=TAUPE, top=2.0, size=18, italic=False):
    box = slide.shapes.add_textbox(Inches(0.85), Inches(top), Inches(11.4), Inches(1.2))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Manrope"
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.italic = italic
    return box


def add_bullets(slide, items, left=0.9, top=2.0, width=11.5, height=4.8):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(8)
        p.font.name = "Manrope"
        p.font.size = Pt(24)
        p.font.color.rgb = INK
    return box


def add_label(slide, text, x=0.85, y=0.25, w=4.2, h=0.4, fg=TAUPE):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Manrope"
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = fg
    return box


def add_footer(slide, text="LSN Consulting · Bible graphique + univers onirique · 2026"):
    box = slide.shapes.add_textbox(Inches(0.85), Inches(7.0), Inches(12.0), Inches(0.3))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.name = "Manrope"
    r.font.size = Pt(10)
    r.font.color.rgb = TAUPE


def color_chip(slide, x, y, color, name, hex_code, role):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.9), Inches(1.6))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.color.rgb = CREAM if color == INK else TAUPE

    t = slide.shapes.add_textbox(Inches(x + 0.12), Inches(y + 1.7), Inches(2.8), Inches(1.1))
    tf = t.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = name
    p1.font.name = "Manrope"
    p1.font.size = Pt(13)
    p1.font.bold = True
    p1.font.color.rgb = INK

    p2 = tf.add_paragraph()
    p2.text = f"{hex_code} · {role}"
    p2.font.name = "Manrope"
    p2.font.size = Pt(10)
    p2.font.color.rgb = TAUPE


# Slide 1: Cover
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, INK)
add_label(s, "LSN CONSULTING · CHARTE GRAPHIQUE", fg=CREAM)
add_title(s, "Elements de charte", color=CREAM, top=1.35, size=58)
add_subtitle(s, "Version recreee depuis le site, avec signature white clay onirique", color=BLUE, top=2.65, size=23, italic=True)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85), Inches(3.45), Inches(2.7), Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
add_subtitle(s, "Source de verite: bible_graphique_lsn_v2.html", color=CREAM, top=4.0, size=16)
add_subtitle(s, "Leila Seghirate-Nigg · Mai 2026", color=CREAM, top=6.6, size=12)

# Slide 2: Intention
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, CREAM)
add_label(s, "01 · INTENTION")
add_title(s, "Un systeme clair, vivant et coherent", size=44)
add_bullets(s, [
    "Aligner toute production visuelle sur la meme identite.",
    "Passer d'une charte theorique a une charte actionnable.",
    "Unifier site, reseaux, supports de vente et presentations."
], top=2.15)
add_footer(s)

# Slide 3: Palette
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
add_label(s, "02 · PALETTE OFFICIELLE")
add_title(s, "4 couleurs fondatrices", size=42)
color_chip(s, 0.85, 2.0, BLUE, "Bleu electrique", "#3333CC", "Accent / CTA")
color_chip(s, 3.95, 2.0, INK, "Anthracite", "#2A2A34", "Titres / fonds sombres")
color_chip(s, 7.05, 2.0, CREAM, "Creme", "#EDE9E3", "Fond / texte sur sombre")
color_chip(s, 10.15, 2.0, TAUPE, "Taupe", "#948F8A", "Secondaire / meta")
add_subtitle(s, "Interdit: ajouter des teintes hors systeme LSN.", top=5.9, size=16)
add_footer(s)

# Slide 4: Typography
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, CREAM)
add_label(s, "03 · TYPOGRAPHIE")
add_title(s, "Noto Serif + Manrope", size=44)
add_subtitle(s, "Noto Serif pour l'empreinte editoriale.", top=2.2, size=18)
add_subtitle(s, "Manrope pour l'interface, la lisibilite et les labels.", top=2.65, size=18)
box = s.shapes.add_textbox(Inches(0.9), Inches(3.35), Inches(12.0), Inches(2.8))
tf = box.text_frame
p = tf.paragraphs[0]
p.text = "Mes realisations"
p.font.name = "Noto Serif"
p.font.bold = True
p.font.size = Pt(56)
p.font.color.rgb = INK
p2 = tf.add_paragraph()
p2.text = "Solutions digitales pour la formation"
p2.font.name = "Manrope"
p2.font.size = Pt(20)
p2.font.bold = True
p2.font.color.rgb = BLUE
add_footer(s)

# Slide 5: White clay universe
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
add_label(s, "04 · UNIVERS GRAPHIQUE")
add_title(s, "Signature: white clay elegant et onirique", size=40)
add_bullets(s, [
    "Matiere dominante: argile blanche mate, douce et sculpturale.",
    "Lumiere: cinematographique diffuse, halos subtils, ombres maitrisees.",
    "Palette image: bleu #3333CC, creme #EDE9E3, ombres ink #2A2A34.",
    "Composition: minimalisme premium, espace negatif, narration poetique."
], top=2.1)
add_footer(s)

# Slide 6: Do/Don't universe
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, CREAM)
add_label(s, "05 · CADRE DE PRODUCTION")
add_title(s, "Ce qu'il faut faire / eviter", size=40)
left = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(2.0), Inches(5.95), Inches(4.3))
left.fill.solid(); left.fill.fore_color.rgb = WHITE; left.line.color.rgb = RGBColor(0x22, 0xC5, 0x5E)
right = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.55), Inches(2.0), Inches(5.95), Inches(4.3))
right.fill.solid(); right.fill.fore_color.rgb = WHITE; right.line.color.rgb = RGBColor(0xEF, 0x44, 0x44)

lb = s.shapes.add_textbox(Inches(1.1), Inches(2.25), Inches(5.3), Inches(3.8)).text_frame
lb.text = "A FAIRE"
lb.paragraphs[0].font.name = "Manrope"; lb.paragraphs[0].font.size = Pt(13); lb.paragraphs[0].font.bold = True; lb.paragraphs[0].font.color.rgb = RGBColor(0x22, 0xC5, 0x5E)
for t in [
    "Formes epurees, volume sculptural.",
    "Ambiance contemplative et premium.",
    "Contraste doux et lisible.",
    "Textures minimes, jamais bruyantes."
]:
    p = lb.add_paragraph(); p.text = t; p.font.name = "Manrope"; p.font.size = Pt(18); p.font.color.rgb = INK

rb = s.shapes.add_textbox(Inches(6.8), Inches(2.25), Inches(5.3), Inches(3.8)).text_frame
rb.text = "A EVITER"
rb.paragraphs[0].font.name = "Manrope"; rb.paragraphs[0].font.size = Pt(13); rb.paragraphs[0].font.bold = True; rb.paragraphs[0].font.color.rgb = RGBColor(0xEF, 0x44, 0x44)
for t in [
    "Effets neon, cyberpunk, trop futuristes.",
    "Glossy plastique ou rendu jeu video.",
    "Saturation excessive et multicolore.",
    "Surcharge d'elements decoratifs."
]:
    p = rb.add_paragraph(); p.text = t; p.font.name = "Manrope"; p.font.size = Pt(18); p.font.color.rgb = INK
add_footer(s)

# Slide 7: Components
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
add_label(s, "06 · COMPOSANTS")
add_title(s, "Langage UI du site", size=42)
add_bullets(s, [
    "Boutons pill arrondis (primary bleu, secondary creme).",
    "Badges uppercase, tracking eleve, contraste fort.",
    "Cards arrondies (24px), bordures fines, ombre discrete.",
    "Navigation sticky creme avec lien actif bleu souligne."
], top=2.2)
add_footer(s)

# Slide 8: Contrasts
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, CREAM)
add_label(s, "07 · ACCESSIBILITE")
add_title(s, "Contrastes a respecter (WCAG AA)", size=40)
add_bullets(s, [
    "Ink sur surface: 15.3:1 (AA)",
    "Bleu sur surface: 7.2:1 (AA)",
    "Creme sur ink: 13.1:1 (AA)",
    "Taupe sur surface: reserve aux textes secondaires/grands corps"
], top=2.2)
add_footer(s)

# Slide 9: Motion
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
add_label(s, "08 · MOTION")
add_title(s, "Mouvement sobre et intentionnel", size=40)
add_bullets(s, [
    "Transitions courtes: 0.2s ease pour l'interface.",
    "Reveals sections: opacity + translateY en 0.5s.",
    "Cards au hover: levee legere, jamais spectaculaire.",
    "Rythme global: calme, premium, maitrise."
], top=2.2)
add_footer(s)

# Slide 10: Prompt ready
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, INK)
add_label(s, "09 · PROMPT MAITRE", fg=CREAM)
add_title(s, "Prompt univers graphique", color=CREAM, size=42)
prompt = (
    '"White clay dreamscape with sculptural characters, LSN electric blue accents (#3333CC), '
    'warm cream atmosphere (#EDE9E3), deep ink shadows (#2A2A34), refined editorial composition, '
    'elegant and poetic mood, soft cinematic light, minimal and premium."'
)
box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(2.2), Inches(11.8), Inches(2.6))
box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x1F, 0x1F, 0x28); box.line.color.rgb = RGBColor(0x44, 0x44, 0x88)
tf = box.text_frame
tf.clear()
p = tf.paragraphs[0]
p.text = prompt
p.font.name = "Noto Serif"
p.font.size = Pt(22)
p.font.italic = True
p.font.color.rgb = CREAM
add_subtitle(s, "A utiliser dans Claude/ChatGPT/Midjourney pour verrouiller l'univers LSN", color=CREAM, top=5.35, size=15)
add_footer(s)

# Slide 11: Mapping old->new
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, CREAM)
add_label(s, "10 · ALIGNEMENT")
add_title(s, "Transformation de l'ancien deck", size=40)
add_bullets(s, [
    "Conserve: intention pedagogique, themes IA, structure de marque.",
    "Remplace: anciennes palettes beige/navy/copper non LSN.",
    "Reformule: vocabulaire visuel vers white clay onirique.",
    "Normalise: composants, typographies, contrastes et ton editorial."
], top=2.2)
add_footer(s)

# Slide 12: Closing
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
add_label(s, "11 · CLOTURE")
add_title(s, "LSN Consulting", size=54)
add_subtitle(s, "Une identite nette, sensible et memorisable.", color=BLUE, top=2.25, size=26, italic=True)
add_bullets(s, [
    "Prochaine etape: decliner ce systeme sur templates reseaux + decks commerciaux.",
    "Source officielle: bible_graphique_lsn_v2.html"
], top=3.5, height=2.3)
add_footer(s)

output = r"C:\Users\seghi\OneDrive\Bureau\LSN FORMATION 2026\CHARTES\bible_graphique\adn charte graphique\ELEMENTS DE CHARTES_LSN_REBUILD_2026.pptx"
prs.save(output)
print(output)
